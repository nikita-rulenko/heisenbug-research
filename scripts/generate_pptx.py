#!/usr/bin/env python3
"""Generate Heisenbug 2026 presentation: AI Agent Context for Testing."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = Path(__file__).parent
TABLES_DIR = BASE / "tables" / "light"
DIAGRAMS_DIR = BASE / "diagrams" / "light_png"
OUT = BASE / "presentation.pptx"

BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
CLR_TITLE = RGBColor(0x1A, 0x1A, 0x2E)
CLR_ACCENT = RGBColor(0x00, 0x6B, 0xA6)
CLR_SUBTITLE = RGBColor(0x55, 0x55, 0x55)
CLR_BODY = RGBColor(0x33, 0x33, 0x33)
CLR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLR_THESIS_BG = RGBColor(0x00, 0x3D, 0x5C)
CLR_LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
CLR_GREEN = RGBColor(0x00, 0x80, 0x60)
CLR_ORANGE = RGBColor(0xD4, 0x6B, 0x08)

W = Inches(13.333)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, color=CLR_BODY, bold=False, alignment=PP_ALIGN.LEFT):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return p


def add_paragraph(tf, text, size=18, color=CLR_BODY, bold=False, alignment=PP_ALIGN.LEFT,
                  space_before=Pt(6), space_after=Pt(2)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


# ── Slide builders ─────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, BG_DARK)

    tf = add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5)).text_frame
    set_text(tf, "Контекст AI-агентов для тестирования", 42, CLR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    tf2 = add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(1)).text_frame
    set_text(tf2, "от хранения фактов к связыванию причин", 28, CLR_ACCENT, alignment=PP_ALIGN.CENTER)

    tf3 = add_textbox(slide, Inches(1), Inches(5), Inches(11), Inches(0.8)).text_frame
    set_text(tf3, "Heisenbug 2026  ·  Никита Руленко", 20, RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.CENTER)


def slide_thesis(prs, text, source="", number=""):
    """Standalone thesis slide — dark background, centered text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CLR_THESIS_BG)

    if number:
        tf_n = add_textbox(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.6)).text_frame
        set_text(tf_n, number, 14, RGBColor(0x66, 0x99, 0xAA), alignment=PP_ALIGN.CENTER)

    tf = add_textbox(slide, Inches(1.2), Inches(2), Inches(10.8), Inches(3)).text_frame
    tf.word_wrap = True
    set_text(tf, text, 32, CLR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    if source:
        tf2 = add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.8)).text_frame
        set_text(tf2, source, 14, RGBColor(0x88, 0xBB, 0xCC), alignment=PP_ALIGN.CENTER)


def slide_section(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CLR_ACCENT)
    tf = add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5)).text_frame
    set_text(tf, title, 40, CLR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if subtitle:
        tf2 = add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(1)).text_frame
        set_text(tf2, subtitle, 22, RGBColor(0xDD, 0xEE, 0xFF), alignment=PP_ALIGN.CENTER)


def slide_content(prs, title, bullets, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_WHITE)

    # accent bar
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), H).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = CLR_ACCENT
    slide.shapes[-1].line.fill.background()

    tf_t = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8)).text_frame
    set_text(tf_t, title, 30, CLR_TITLE, bold=True)

    tf = add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5)).text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            set_text(tf, f"  {b}", 20, CLR_BODY)
        else:
            add_paragraph(tf, f"  {b}", 20, CLR_BODY)

    if note:
        tf_n = add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.7)).text_frame
        set_text(tf_n, note, 13, CLR_SUBTITLE)


def slide_image(prs, title, img_path, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_WHITE)

    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), H).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = CLR_ACCENT
    slide.shapes[-1].line.fill.background()

    tf_t = add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.7)).text_frame
    set_text(tf_t, title, 26, CLR_TITLE, bold=True)

    if subtitle:
        tf_s = add_textbox(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.5)).text_frame
        set_text(tf_s, subtitle, 14, CLR_SUBTITLE)

    img = str(img_path)
    max_w = Inches(12)
    max_h = Inches(5.8)
    top_offset = Inches(1.3)

    from PIL import Image as PILImage
    with PILImage.open(img) as im:
        iw, ih = im.size
    aspect = iw / ih
    if max_w / max_h > aspect:
        h = max_h
        w = int(h * aspect)
    else:
        w = max_w
        h = int(w / aspect)
    left = int((W - w) / 2)
    slide.shapes.add_picture(img, left, int(top_offset), w, h)


def slide_two_column(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_WHITE)

    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), H).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = CLR_ACCENT
    slide.shapes[-1].line.fill.background()

    tf_t = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8)).text_frame
    set_text(tf_t, title, 30, CLR_TITLE, bold=True)

    # left column
    tf_l_title = add_textbox(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(0.6)).text_frame
    set_text(tf_l_title, left_title, 22, CLR_ACCENT, bold=True)

    tf_l = add_textbox(slide, Inches(0.8), Inches(1.9), Inches(5.6), Inches(5)).text_frame
    tf_l.word_wrap = True
    for i, b in enumerate(left_bullets):
        if i == 0:
            set_text(tf_l, f"  {b}", 17, CLR_BODY)
        else:
            add_paragraph(tf_l, f"  {b}", 17, CLR_BODY)

    # right column
    tf_r_title = add_textbox(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(0.6)).text_frame
    set_text(tf_r_title, right_title, 22, CLR_GREEN, bold=True)

    tf_r = add_textbox(slide, Inches(7.0), Inches(1.9), Inches(5.6), Inches(5)).text_frame
    tf_r.word_wrap = True
    for i, b in enumerate(right_bullets):
        if i == 0:
            set_text(tf_r, f"  {b}", 17, CLR_BODY)
        else:
            add_paragraph(tf_r, f"  {b}", 17, CLR_BODY)


def slide_end(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)

    tf = add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(1.5)).text_frame
    set_text(tf, "Спасибо!", 48, CLR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    tf2 = add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(1)).text_frame
    set_text(tf2, "Как должна выглядеть AI-память для тестирования через 2 года?",
             24, CLR_ACCENT, alignment=PP_ALIGN.CENTER)

    tf3 = add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(1.2)).text_frame
    set_text(tf3, "github.com/nikita-rulenko/helixir", 18,
             RGBColor(0x88, 0xBB, 0xCC), alignment=PP_ALIGN.CENTER)
    add_paragraph(tf3, "Никита Руленко  ·  Heisenbug 2026", 16,
                  RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.CENTER, space_before=Pt(12))


# ── Main ───────────────────────────────────────────────────────

def main():
    prs = new_prs()

    # ═══════════════════════════════════════════════════════════
    # 1. TITLE
    # ═══════════════════════════════════════════════════════════
    slide_title(prs)

    # ═══════════════════════════════════════════════════════════
    # 2. PROBLEM
    # ═══════════════════════════════════════════════════════════
    slide_section(prs, "Проблема", "Почему context window — не ответ")

    slide_thesis(prs,
                 "Тестов всегда больше\nчем production-кода",
                 "Miranda et al., 2025: test code = 1:1 – 1:3 от production code\n"
                 "Mature backend: 50-55% кодовой базы — тесты",
                 number="ПРОБЛЕМА")

    slide_content(prs, "6 проблем контекста в тестировании", [
        "① Context window exhaustion — 500 тест-файлов не влезут в 128K tokens",
        "② Context rot — 15 загруженных файлов = деградация генерации",
        "③ Pattern drift — конвенции меняются: моки → testcontainers",
        "④ Knowledge loss — почему тест написан так? Между сессиями теряется",
        "⑤ Flaky test debugging — race conditions переоткрываются каждый раз",
        "⑥ Test maintenance debt — UI меняется → десятки тестов ломаются",
    ], "OAJAIML, Jan 2026: MECW значительно меньше заявленного · "
       "Lost in the Middle: LLM performance -20% для информации в середине контекста")

    slide_thesis(prs,
                 "Увеличение context window\nне решит проблему",
                 "OAJAIML, 2026: \"Large context windows degrade model performance.\n"
                 "Agentic systems relying on large context will see cascading failures.\"\n\n"
                 "Stephanie Jarmak, 2026: \"Enterprise = 100B LOC. Context window = 200K max effective (The Hobbit).\"",
                 number="ТЕЗИС")

    slide_thesis(prs,
                 "Те, кто получает максимум от coding agents —\n"
                 "те, у кого сильные тестовые практики",
                 "Addy Osmani, 2026: \"An agent can 'fly' through a project\n"
                 "with a good test suite as safety net. Without tests,\n"
                 "the agent might blithely assume everything is fine.\"",
                 number="ТЕЗИС")

    # ═══════════════════════════════════════════════════════════
    # 3. LAYER 1: MD FILES
    # ═══════════════════════════════════════════════════════════
    slide_section(prs, "Слой 1: MD файлы", "Policy layer — конституция для агента")

    slide_content(prs, "MD файлы: как это работает", [
        "AGENTS.md / .cursorrules / spec.md — статичные файлы в репозитории",
        "Агент читает при старте каждой сессии",
        "Содержит: test framework, команды запуска, конвенции, запреты",
        "60K+ проектов используют AGENTS.md (Linux Foundation, 2026)",
        "Version-controlled, zero cost, zero latency",
        "",
        "✅ Мгновенный onboarding нового агента",
        "✅ Понятно человеку — можно review-ить",
        "❌ Статичен — не учится на ошибках",
        "❌ Не масштабируется при 500+ тестах",
    ])

    slide_image(prs, "MD файлы — flow контекста",
                DIAGRAMS_DIR / "flow_md_files.png",
                "Однонаправленный поток: файл → prompt → агент")

    slide_thesis(prs,
                 "MD файлы — конституция для агента.\nНо конституция не учится.",
                 "При росте до 500+ тестов — недостаточно.\n"
                 "Augment Code, 2026: \"rules not set to always apply were easy for the agent to miss\"",
                 number="ТЕЗИС")

    # ═══════════════════════════════════════════════════════════
    # 4. LAYER 2: GITHUB ISSUES MCP
    # ═══════════════════════════════════════════════════════════
    slide_section(prs, "Слой 2: GitHub Issues MCP", "Coordination layer — кто что делает")

    slide_content(prs, "GitHub Issues MCP: как это работает", [
        "CI failure → GitHub Issue → AI Agent (MCP) → fix → PR → close",
        "Bidirectional audit trail: QA и AI в одном пространстве",
        "GitHub Copilot Coding Agent: assign issue → 👀 → draft PR → comments",
        "Free tier: 50 premium requests/month — достаточно для старта",
        "",
        "✅ Full audit trail каждого бага",
        "✅ Collaboration: комментарии, labels, assignees",
        "✅ CI integration: auto-issue → auto-fix → auto-close",
        "❌ Нет semantic memory — каждый баг с нуля",
        "❌ Нет обобщения паттернов",
    ], "GitHub Agentic Workflows, 2026: \"60% of investigations lead to fixes within 24-48 hours\"")

    slide_image(prs, "GitHub Issues MCP — flow контекста",
                DIAGRAMS_DIR / "flow_github_mcp.png",
                "Bidirectional: issue ↔ agent ↔ PR")

    slide_thesis(prs,
                 "Issues — это workflow tool, не memory.\n"
                 "Используйте для трекинга задач,\nне ожидайте что агент «помнит».",
                 "",
                 number="ТЕЗИС")

    # ═══════════════════════════════════════════════════════════
    # 5. LAYER 3: MEM0
    # ═══════════════════════════════════════════════════════════
    slide_section(prs, "Слой 3: Mem0", "Knowledge layer — production-ready AI-память")

    slide_content(prs, "Mem0: как это работает", [
        "Агент ведёт беседу → Mem0 автоматически извлекает факты (LLM)",
        "Факты хранятся как embeddings в vector DB + entity graph",
        "При следующей сессии — semantic search по смыслу",
        "A.U.D.N. цикл: ADD / UPDATE / DELETE / NOOP — LLM решает",
        "47.8K GitHub stars, $24M funding, SOC 2 Type II",
        "Интеграции: CrewAI, LangChain, AWS Agent SDK, Vercel AI",
        "",
        "✅ Cross-session memory — агент помнит между сессиями",
        "✅ Auto-extraction — не нужно вручную обновлять .md",
        "✅ Cloud Free: 10K memories — достаточно для старта",
    ])

    slide_image(prs, "Mem0 — архитектура и flow",
                DIAGRAMS_DIR / "flow_mem0_wide.png",
                "LLM extraction → Vector + Graph DB → Semantic retrieval")

    slide_thesis(prs,
                 "90% сокращение токенов\n"
                 "26% рост accuracy\n"
                 "91% снижение latency",
                 "Mem0 LOCOMO benchmark, 2025\n"
                 "vs full-context подход",
                 number="BENCHMARK")

    slide_thesis(prs,
                 "При противоречии Mem0 удаляет\nстарый факт. История потеряна.",
                 "A.U.D.N. цикл: DELETE = старая memory удаляется навсегда.\n"
                 "UPDATE = старое значение перезаписывается.\n"
                 "Нет аудита: \"что было до того, как LLM решил обновить?\"",
                 number="ОГРАНИЧЕНИЕ")

    slide_two_column(prs, "Mem0: Platform (Cloud) vs Open Source",
                     "Cloud (Managed)", [
                         "⏱ 5 минут до первой memory",
                         "💰 Free: 10K memories",
                         "💰 Starter: $19/mo, Pro: $249/mo",
                         "✅ SOC 2, GDPR из коробки",
                         "✅ Auto-scaling, HA",
                         "✅ Webhooks, export, filters v2",
                         "❌ Данные на серверах Mem0 (US)",
                     ],
                     "Self-hosted (OSS)", [
                         "⏱ 15-30 мин → часы на настройку",
                         "💰 $0 + инфраструктура + LLM tokens",
                         "🔧 Vector DB + Graph DB + LLM = 3 сервиса",
                         "✅ Полный контроль над данными",
                         "✅ Любая юрисдикция",
                         "❌ Manual scaling, нет HA",
                         "❌ DevOps overhead",
                     ])

    # ═══════════════════════════════════════════════════════════
    # 6. LAYER 4: HELIXIR
    # ═══════════════════════════════════════════════════════════
    slide_section(prs, "Слой 4: Helixir", "Reasoning layer — proof of concept")

    slide_thesis(prs,
                 "Хранение фактов ≠ связывание фактов",
                 "Mem0 отвечает: «что мы знаем?»\n"
                 "Каузальная память отвечает: «почему мы это знаем и к чему это ведёт?»\n\n"
                 "Готового production-решения для связывания — нет.",
                 number="КЛЮЧЕВОЙ ТЕЗИС")

    slide_content(prs, "Helixir: что он показывает", [
        "Первый каузальный фреймворк AI-памяти (proof of concept)",
        "",
        "Каузальные цепочки: IMPLIES / BECAUSE / CONTRADICTS / SUPERSEDE",
        "Онтология: 7 типов (skill, preference, goal, fact, opinion, experience, achievement)",
        "FastThink: изолированный scratchpad — промежуточные мысли не в памяти",
        "Temporal filtering: 4 режима (4ч / 30д / 90д / all)",
        "Append-only graph: ничего не удаляется, полная история",
        "",
        "⚠️  v0.1.1, 74 stars, 1 maintainer — это НЕ production",
        "⚠️  LLM обязателен (как и Mem0) — extraction + dedup",
        "⚠️  Нет benchmarks, нет production case studies",
    ])

    slide_image(prs, "Helixir — архитектура и flow",
                DIAGRAMS_DIR / "flow_helixir_wide.png",
                "LLM extraction → Causal graph (HelixDB) → SmartTraversal")

    slide_two_column(prs, "Mem0 vs Helixir: при противоречии",
                     "Mem0: DELETE/UPDATE", [
                         "Тест падал по 3 причинам:",
                         "  Jan: timeout",
                         "  Feb: race condition → UPDATE → timeout удалён",
                         "  Mar: cookie settings → UPDATE → race удалён",
                         "",
                         "В памяти: только «cookie settings»",
                         "История timeout и race condition — потеряна",
                         "При 4-м падении: агент не знает",
                         "о предыдущих причинах",
                     ],
                     "Helixir: SUPERSEDE", [
                         "Тест падал по 3 причинам:",
                         "  Jan: timeout (fact)",
                         "  Feb: race condition → SUPERSEDE → timeout",
                         "  Mar: cookie settings → SUPERSEDE → race",
                         "",
                         "В памяти: ВСЕ 3 причины + цепочка",
                         "timeout → race → cookie (эволюция)",
                         "При 4-м падении: агент проверяет",
                         "все 3 предыдущие причины",
                     ])

    slide_thesis(prs,
                 "Helixir — первый пробный шар.\n"
                 "Proof of concept, не продукт.",
                 "1 maintainer — это несерьёзно для production зависимости.\n"
                 "Но это направление, в которое может прийти AI-память.\n\n"
                 "Сегодня мы храним факты. Завтра будем хранить связи между фактами.",
                 number="ПОЗИЦИОНИРОВАНИЕ")

    # ═══════════════════════════════════════════════════════════
    # 7. COMPARISON
    # ═══════════════════════════════════════════════════════════
    slide_section(prs, "Сравнение", "4 слоя, а не 4 альтернативы")

    slide_image(prs, "Сравнительная таблица: 4 подхода",
                TABLES_DIR / "10_slide_comparison.png",
                "✓ = есть  ✗ = нет  ★ = лучший в критерии  ~ = частично")

    slide_thesis(prs,
                 "Четыре слоя, а не четыре альтернативы",
                 "MD (policy) → Issues (coordination) → Mem0 (knowledge) → каузальная память (reasoning)\n\n"
                 "Не «выберите один», а «добавляйте следующий по мере роста боли»",
                 number="ГЛАВНАЯ МЫСЛЬ")

    # ═══════════════════════════════════════════════════════════
    # 8. TESTING SCENARIOS
    # ═══════════════════════════════════════════════════════════
    slide_section(prs, "Сценарии", "Как это работает на практике")

    slide_content(prs, "Сценарий: Flaky тест падает в CI", [
        "MD файлы: в AGENTS.md может быть «check for race conditions first»",
        "  → Общая рекомендация, не конкретная",
        "",
        "GitHub Issues: CI failure → issue → agent investigates → PR",
        "  → Audit trail есть, но каждый раз с нуля",
        "",
        "Mem0: «в прошлый раз login test падал, фиксили увеличением timeout»",
        "  → Помнит факт, но не причинно-следственную связь",
        "",
        "Helixir: «Login test fails BECAUSE session race condition (3 раза).",
        "  Fix: add retry + увеличение timeout до 30s»",
        "  → Каузальная цепочка: причина → следствие → решение",
    ])

    slide_content(prs, "Сценарий: Эволюция тест-стратегии", [
        "MD файлы: обновляются вручную, могут отставать от реальности",
        "",
        "GitHub Issues: нет механизма трекинга эволюции",
        "",
        "Mem0: помнит текущие практики, но не связывает «было → стало»",
        "  → «Используем Playwright» — а что было до этого?",
        "",
        "Helixir: явная эволюция через SUPERSEDE",
        "  «Selenium» SUPERSEDE→ «Playwright»",
        "  «Mock DB» SUPERSEDE→ «Testcontainers»",
        "  → Агент не предложит устаревший паттерн",
    ])

    # ═══════════════════════════════════════════════════════════
    # 9. ACTIONABLE TAKEAWAYS
    # ═══════════════════════════════════════════════════════════
    slide_section(prs, "Что делать в понедельник", "3 из 4 слоёв — бесплатны и доступны сегодня")

    slide_content(prs, "Начните в понедельник", [
        "① AGENTS.md с тестовыми конвенциями = 10 минут",
        "   pytest, структура, fixtures, запреты — агент сразу знает правила",
        "",
        "② GitHub MCP в Cursor = 5 минут",
        "   CI failure → issue → agent → PR — audit trail из коробки",
        "",
        "③ Mem0 Cloud Free = 5 минут",
        "   10K memories бесплатно — агент помнит между сессиями",
        "",
        "④ Каузальная память = ?",
        "   Production-решения пока нет. Helixir — proof of concept.",
        "   Следите за развитием: graph memory + causal chains",
    ])

    # ═══════════════════════════════════════════════════════════
    # 10. CLOSING
    # ═══════════════════════════════════════════════════════════
    slide_thesis(prs,
                 "Это открытый вопрос.\n"
                 "Доклад не закрывает тему —\nдоклад открывает дискуссию.",
                 "Как должна выглядеть AI-память для тестирования через 2 года?\n\n"
                 "Хранение фактов — решено (Mem0).\n"
                 "Связывание фактов — открытая задача.",
                 number="ФИНАЛ")

    slide_end(prs)

    prs.save(str(OUT))
    print(f"✅ Saved: {OUT}")
    print(f"   Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
